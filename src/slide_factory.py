"""
スライド種別ごとの生成関数
カラースキームのキーを参照して色を決定する
"""
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData

from src.style import apply_font, hex_to_rgb, LAYOUT


def _get_color(style, key, fallback):
    """スタイルからカラーを取得するヘルパー"""
    color_hex = style.get(key, fallback)
    return hex_to_rgb(color_hex)


def add_title_slide(prs, data, style):
    """タイトルスライドを追加"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白レイアウト

    accent = _get_color(style, "accent_color", "4472C4")
    text_color = _get_color(style, "text_color", "333333")
    heading_color = _get_color(style, "heading_color", "333333")
    bg_color = style.get("bg_color", "FFFFFF")

    # 背景色の適用
    if bg_color != "FFFFFF":
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = hex_to_rgb(bg_color)

    # タイトルテキストボックス
    left = Inches(1.0)
    top = Inches(2.5)
    width = Inches(11.3)
    height = Inches(1.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = data.get("title", "")
    apply_font(run, style.get("font_name", "Meiryo"),
               style.get("title_size_pt", 28), heading_color, bold=True)

    # サブタイトル
    subtitle = data.get("subtitle", "")
    if subtitle:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(12)
        run2 = p2.add_run()
        run2.text = subtitle
        apply_font(run2, style.get("font_name", "Meiryo"),
                   style.get("subtitle_size_pt", 18), accent)

    # 日付・作成者
    author = data.get("author", "")
    date = data.get("date", "")
    if author or date:
        info_text = " | ".join(filter(None, [author, date]))
        p3 = tf.add_paragraph()
        p3.alignment = PP_ALIGN.CENTER
        p3.space_before = Pt(24)
        run3 = p3.add_run()
        run3.text = info_text
        apply_font(run3, style.get("font_name", "Meiryo"), 12, text_color)

    # アクセントライン
    line_left = Inches(4.0)
    line_top = Inches(2.2)
    line_width = Inches(5.3)
    line = slide.shapes.add_shape(
        1, line_left, line_top, line_width, Pt(4)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = accent
    line.line.fill.background()


def add_content_slide(prs, data, style):
    """箇条書きコンテンツスライドを追加"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    accent = _get_color(style, "accent_color", "4472C4")
    text_color = _get_color(style, "text_color", "333333")
    bg_color = style.get("bg_color", "FFFFFF")

    # 背景色
    if bg_color != "FFFFFF":
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = hex_to_rgb(bg_color)

    # スライドタイトル
    _add_slide_title(slide, data.get("title", ""), style, accent)

    # 本文（箇条書き）
    body = data.get("body", [])
    if isinstance(body, str):
        body = [body]

    left = LAYOUT["margin_left"]
    top = LAYOUT["margin_top"]
    width = LAYOUT["content_width"]
    height = LAYOUT["content_height"]
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    font_name = style.get("font_name", "Meiryo")
    body_size = style.get("body_size_pt", 14)

    for i, item in enumerate(body):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(8)
        p.level = 0

        # 箇条書きマーカー
        run = p.add_run()
        run.text = f"● {item}"
        apply_font(run, font_name, body_size, text_color)


def add_table_slide(prs, data, style):
    """テーブルスライドを追加"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    accent = _get_color(style, "accent_color", "4472C4")
    text_color = _get_color(style, "text_color", "333333")
    bg_color = style.get("bg_color", "FFFFFF")

    # テーブル固有の色（カラースキームから取得）
    header_bg = _get_color(style, "table_header_bg", style.get("accent_color", "4472C4"))
    header_text = _get_color(style, "table_header_text", "FFFFFF")
    alt_row = _get_color(style, "table_alt_row", "F2F2F2")

    # 背景色
    if bg_color != "FFFFFF":
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = hex_to_rgb(bg_color)

    _add_slide_title(slide, data.get("title", ""), style, accent)

    headers = data.get("headers", [])
    rows = data.get("rows", [])
    if not headers and not rows:
        return

    num_rows = len(rows) + 1
    num_cols = len(headers) if headers else (len(rows[0]) if rows else 0)
    if num_cols == 0:
        return

    left = LAYOUT["margin_left"]
    top = LAYOUT["margin_top"]
    width = LAYOUT["content_width"]
    height = Inches(0.4) * num_rows

    table_shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, height)
    table = table_shape.table

    font_name = style.get("font_name", "Meiryo")

    # ヘッダー行
    for j, header in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = str(header)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_bg
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                apply_font(run, font_name, 12, header_text, bold=True)

    # データ行
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            if j >= num_cols:
                break
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            # 交互の行色
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = alt_row
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    apply_font(run, font_name, 11, text_color)


def add_chart_slide(prs, data, style):
    """グラフスライドを追加"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    accent = _get_color(style, "accent_color", "4472C4")
    bg_color = style.get("bg_color", "FFFFFF")

    # 背景色
    if bg_color != "FFFFFF":
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = hex_to_rgb(bg_color)

    _add_slide_title(slide, data.get("title", ""), style, accent)

    chart_type_str = data.get("chart_type", "bar")
    categories = data.get("categories", [])
    series_list = data.get("series", [])

    if not categories or not series_list:
        return

    # チャートタイプのマッピング
    chart_type_map = {
        "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "line": XL_CHART_TYPE.LINE,
        "pie": XL_CHART_TYPE.PIE,
    }
    chart_type = chart_type_map.get(chart_type_str, XL_CHART_TYPE.COLUMN_CLUSTERED)

    chart_data = CategoryChartData()
    chart_data.categories = categories
    for s in series_list:
        chart_data.add_series(s.get("name", ""), s.get("values", []))

    left = LAYOUT["margin_left"]
    top = LAYOUT["margin_top"]
    width = LAYOUT["content_width"]
    height = LAYOUT["content_height"]

    chart_shape = slide.shapes.add_chart(chart_type, left, top, width, height, chart_data)

    chart = chart_shape.chart
    chart.has_legend = len(series_list) > 1
    if chart.has_legend:
        chart.legend.include_in_layout = False


def add_image_slide(prs, data, style):
    """画像スライドを追加"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    accent = _get_color(style, "accent_color", "4472C4")
    text_color = _get_color(style, "text_color", "333333")
    bg_color = style.get("bg_color", "FFFFFF")

    # 背景色
    if bg_color != "FFFFFF":
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = hex_to_rgb(bg_color)

    _add_slide_title(slide, data.get("title", ""), style, accent)

    image_path = data.get("image_path", "")
    if not image_path:
        return

    try:
        left = Inches(2.0)
        top = LAYOUT["margin_top"]
        height = LAYOUT["content_height"]
        slide.shapes.add_picture(image_path, left, top, height=height)
    except Exception:
        txBox = slide.shapes.add_textbox(
            LAYOUT["margin_left"], LAYOUT["margin_top"],
            LAYOUT["content_width"], Inches(1.0)
        )
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = f"[画像が見つかりません: {image_path}]"
        apply_font(run, style.get("font_name", "Meiryo"), 14, text_color)


def add_section_slide(prs, data, style):
    """セクション区切りスライドを追加"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # セクション固有の色（カラースキームから取得）
    section_bg = _get_color(style, "section_bg", style.get("accent_color", "4472C4"))
    section_text = _get_color(style, "section_text", "FFFFFF")

    # 背景塗りつぶし
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = section_bg

    # セクションタイトル（中央配置）
    left = Inches(1.0)
    top = Inches(2.8)
    width = Inches(11.3)
    height = Inches(1.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = data.get("title", "")
    apply_font(run, style.get("font_name", "Meiryo"),
               style.get("heading_size_pt", 24), section_text, bold=True)


def add_two_column_slide(prs, data, style):
    """2カラムレイアウトスライドを追加"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    accent = _get_color(style, "accent_color", "4472C4")
    text_color = _get_color(style, "text_color", "333333")
    bg_color = style.get("bg_color", "FFFFFF")

    if bg_color != "FFFFFF":
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = hex_to_rgb(bg_color)

    _add_slide_title(slide, data.get("title", ""), style, accent)

    font_name = style.get("font_name", "Meiryo")
    body_size = style.get("body_size_pt", 14)
    sub_heading_size = style.get("heading_size_pt", 24) - 6

    col_width = Inches(5.5)
    top = LAYOUT["margin_top"]
    height = LAYOUT["content_height"]
    left_x = LAYOUT["margin_left"]
    right_x = left_x + col_width + Inches(0.7)

    # 左カラム
    _render_column(slide, left_x, top, col_width, height,
                   data.get("left_title", ""), data.get("left_body", []),
                   font_name, sub_heading_size, body_size, accent, text_color)

    # 右カラム
    _render_column(slide, right_x, top, col_width, height,
                   data.get("right_title", ""), data.get("right_body", []),
                   font_name, sub_heading_size, body_size, accent, text_color)

    # 中央の区切り線
    divider_x = left_x + col_width + Inches(0.3)
    line = slide.shapes.add_shape(
        1, divider_x, top, Pt(1.5), height
    )
    line.fill.solid()
    line.fill.fore_color.rgb = _get_color(style, "accent_light", style.get("accent_color", "4472C4"))
    line.line.fill.background()


def add_key_message_slide(prs, data, style):
    """キーメッセージスライドを追加（大きいメッセージ + 根拠）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    accent = _get_color(style, "accent_color", "4472C4")
    text_color = _get_color(style, "text_color", "333333")
    bg_color = style.get("bg_color", "FFFFFF")

    if bg_color != "FFFFFF":
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = hex_to_rgb(bg_color)

    _add_slide_title(slide, data.get("title", ""), style, accent)

    font_name = style.get("font_name", "Meiryo")
    body_size = style.get("body_size_pt", 14)

    # キーメッセージ（大きく中央寄せ）
    message = data.get("message", "")
    if message:
        msg_top = LAYOUT["margin_top"]
        txBox = slide.shapes.add_textbox(
            Inches(1.0), msg_top, Inches(11.3), Inches(1.8)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = message
        apply_font(run, font_name, 26, accent, bold=True)

        # メッセージ下のアクセントライン
        line_top = msg_top + Inches(2.0)
        line = slide.shapes.add_shape(
            1, Inches(3.0), line_top, Inches(7.3), Pt(2)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = _get_color(style, "accent_light", style.get("accent_color", "4472C4"))
        line.line.fill.background()

    # 根拠・サポート箇条書き
    body = data.get("body", [])
    if isinstance(body, str):
        body = [body]

    if body:
        body_top = LAYOUT["margin_top"] + Inches(2.5)
        txBox = slide.shapes.add_textbox(
            Inches(1.5), body_top, Inches(10.3), Inches(3.0)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, item in enumerate(body):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.space_after = Pt(6)
            run = p.add_run()
            run.text = f"● {item}"
            apply_font(run, font_name, body_size, text_color)


def add_comparison_slide(prs, data, style):
    """Before/After比較スライドを追加"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    accent = _get_color(style, "accent_color", "4472C4")
    text_color = _get_color(style, "text_color", "333333")
    bg_color = style.get("bg_color", "FFFFFF")

    if bg_color != "FFFFFF":
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = hex_to_rgb(bg_color)

    _add_slide_title(slide, data.get("title", ""), style, accent)

    font_name = style.get("font_name", "Meiryo")
    body_size = style.get("body_size_pt", 14)

    col_width = Inches(5.2)
    top = LAYOUT["margin_top"]
    height = LAYOUT["content_height"]
    left_x = LAYOUT["margin_left"]
    right_x = left_x + col_width + Inches(1.3)

    # Beforeラベル背景（赤系）
    before_label_bg = slide.shapes.add_shape(
        1, left_x, top, col_width, Inches(0.5)
    )
    before_label_bg.fill.solid()
    before_label_bg.fill.fore_color.rgb = hex_to_rgb("E74C3C")
    before_label_bg.line.fill.background()
    # Beforeラベルテキスト
    txBox = slide.shapes.add_textbox(left_x, top, col_width, Inches(0.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = data.get("before_title", "Before")
    apply_font(run, font_name, 16, hex_to_rgb("FFFFFF"), bold=True)

    # Before本文
    before_items = data.get("before_items", [])
    if isinstance(before_items, str):
        before_items = [before_items]
    if before_items:
        txBox = slide.shapes.add_textbox(
            left_x + Inches(0.2), top + Inches(0.7),
            col_width - Inches(0.4), height - Inches(0.8)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, item in enumerate(before_items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.space_after = Pt(6)
            run = p.add_run()
            run.text = f"● {item}"
            apply_font(run, font_name, body_size, text_color)

    # 中央の矢印（→）
    arrow_x = left_x + col_width + Inches(0.2)
    arrow_top = top + Inches(1.5)
    txBox = slide.shapes.add_textbox(arrow_x, arrow_top, Inches(0.9), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "→"
    apply_font(run, font_name, 32, accent, bold=True)

    # Afterラベル背景（緑系）
    after_label_bg = slide.shapes.add_shape(
        1, right_x, top, col_width, Inches(0.5)
    )
    after_label_bg.fill.solid()
    after_label_bg.fill.fore_color.rgb = hex_to_rgb("27AE60")
    after_label_bg.line.fill.background()
    # Afterラベルテキスト
    txBox = slide.shapes.add_textbox(right_x, top, col_width, Inches(0.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = data.get("after_title", "After")
    apply_font(run, font_name, 16, hex_to_rgb("FFFFFF"), bold=True)

    # After本文
    after_items = data.get("after_items", [])
    if isinstance(after_items, str):
        after_items = [after_items]
    if after_items:
        txBox = slide.shapes.add_textbox(
            right_x + Inches(0.2), top + Inches(0.7),
            col_width - Inches(0.4), height - Inches(0.8)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, item in enumerate(after_items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.space_after = Pt(6)
            run = p.add_run()
            run.text = f"● {item}"
            apply_font(run, font_name, body_size, text_color)


def _render_column(slide, x, top, width, height, title, body, font_name, title_size, body_size, accent, text_color):
    """2カラムスライドの1カラムを描画するヘルパー"""
    body_top = top
    if title:
        txBox = slide.shapes.add_textbox(x, top, width, Inches(0.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = title
        apply_font(run, font_name, title_size, accent, bold=True)
        body_top = top + Inches(0.6)

    if isinstance(body, str):
        body = [body]

    if body:
        txBox = slide.shapes.add_textbox(x, body_top, width, height - Inches(0.6))
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, item in enumerate(body):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.space_after = Pt(8)
            run = p.add_run()
            run.text = f"● {item}"
            apply_font(run, font_name, body_size, text_color)


def _add_slide_title(slide, title_text, style, accent_color):
    """各スライドのタイトルを追加するヘルパー"""
    if not title_text:
        return

    heading_color = _get_color(style, "heading_color", "333333")

    left = LAYOUT["margin_left"]
    top = LAYOUT["title_top"]
    width = LAYOUT["content_width"]
    height = LAYOUT["title_height"]

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title_text
    apply_font(run, style.get("font_name", "Meiryo"),
               style.get("heading_size_pt", 24), heading_color, bold=True)

    # タイトル下のアクセントライン
    line = slide.shapes.add_shape(
        1, left, top + height - Pt(4), Inches(3.0), Pt(3)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = accent_color
    line.line.fill.background()


# スライドタイプ名と関数のマッピング
SLIDE_TYPES = {
    "title": add_title_slide,
    "content": add_content_slide,
    "two_column": add_two_column_slide,
    "key_message": add_key_message_slide,
    "comparison": add_comparison_slide,
    "table": add_table_slide,
    "chart": add_chart_slide,
    "image": add_image_slide,
    "section": add_section_slide,
}
